#!/usr/bin/env python3
"""
GARAMe Manager 매뉴얼 변환 스크립트
Markdown → DOCX, PPTX 변환

필요 라이브러리:
    pip3 install python-docx python-pptx markdown
"""

import re
import os
from pathlib import Path


def install_dependencies():
    """필요한 라이브러리 설치"""
    import subprocess
    import sys

    required = ['python-docx', 'python-pptx', 'markdown']

    print("필요한 라이브러리 확인 중...")
    for pkg in required:
        try:
            if pkg == 'python-docx':
                import docx
            elif pkg == 'python-pptx':
                import pptx
            elif pkg == 'markdown':
                import markdown
            print(f"✓ {pkg} 설치됨")
        except ImportError:
            print(f"✗ {pkg} 미설치 - 설치 중...")
            subprocess.check_call([sys.executable, "-m", "pip", "install", pkg])
            print(f"✓ {pkg} 설치 완료")


def parse_markdown(md_file):
    """Markdown 파일 파싱"""
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()

    # 구조 분석
    sections = []
    current_section = None
    current_content = []

    for line in content.split('\n'):
        # 제목 감지
        h1_match = re.match(r'^# (.+)$', line)
        h2_match = re.match(r'^## (.+)$', line)
        h3_match = re.match(r'^### (.+)$', line)
        h4_match = re.match(r'^#### (.+)$', line)

        if h1_match:
            if current_section:
                current_section['content'] = '\n'.join(current_content)
                sections.append(current_section)
            current_section = {'level': 1, 'title': h1_match.group(1), 'content': ''}
            current_content = []
        elif h2_match:
            if current_section:
                current_section['content'] = '\n'.join(current_content)
                sections.append(current_section)
            current_section = {'level': 2, 'title': h2_match.group(1), 'content': ''}
            current_content = []
        elif h3_match:
            if current_section:
                current_section['content'] = '\n'.join(current_content)
                sections.append(current_section)
            current_section = {'level': 3, 'title': h3_match.group(1), 'content': ''}
            current_content = []
        elif h4_match:
            if current_section:
                current_section['content'] = '\n'.join(current_content)
                sections.append(current_section)
            current_section = {'level': 4, 'title': h4_match.group(1), 'content': ''}
            current_content = []
        else:
            current_content.append(line)

    # 마지막 섹션 추가
    if current_section:
        current_section['content'] = '\n'.join(current_content)
        sections.append(current_section)

    return sections


def markdown_to_docx(md_file, output_file):
    """Markdown → DOCX 변환"""
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    print(f"\n📄 DOCX 변환 시작: {md_file}")

    doc = Document()

    # 스타일 설정
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Malgun Gothic'  # 맑은 고딕
    font.size = Pt(11)

    # 매뉴얼 파싱
    sections = parse_markdown(md_file)

    for section in sections:
        level = section['level']
        title = section['title']
        content = section['content'].strip()

        # 제목 추가
        if level == 1:
            heading = doc.add_heading(title, level=1)
            heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = heading.runs[0]
            run.font.color.rgb = RGBColor(44, 62, 80)  # #2C3E50
            run.font.size = Pt(24)
            run.font.bold = True
        elif level == 2:
            heading = doc.add_heading(title, level=2)
            run = heading.runs[0]
            run.font.color.rgb = RGBColor(52, 73, 94)  # #34495E
            run.font.size = Pt(18)
            run.font.bold = True
        elif level == 3:
            heading = doc.add_heading(title, level=3)
            run = heading.runs[0]
            run.font.size = Pt(14)
            run.font.bold = True
        elif level == 4:
            heading = doc.add_heading(title, level=4)
            run = heading.runs[0]
            run.font.size = Pt(12)
            run.font.bold = True

        # 내용 추가
        if content:
            # 코드 블록 처리
            code_block_pattern = r'```[\w]*\n(.*?)\n```'
            parts = re.split(code_block_pattern, content, flags=re.DOTALL)

            for i, part in enumerate(parts):
                if i % 2 == 0:
                    # 일반 텍스트
                    for line in part.split('\n'):
                        if line.strip():
                            # 리스트 항목 처리
                            if re.match(r'^[-*]\s+', line):
                                p = doc.add_paragraph(line.strip()[2:], style='List Bullet')
                            elif re.match(r'^\d+\.\s+', line):
                                p = doc.add_paragraph(re.sub(r'^\d+\.\s+', '', line), style='List Number')
                            else:
                                p = doc.add_paragraph(line)

                            # 굵은 글씨 처리
                            if '**' in line:
                                p.clear()
                                parts_bold = re.split(r'\*\*(.*?)\*\*', line)
                                for j, part_text in enumerate(parts_bold):
                                    run = p.add_run(part_text)
                                    if j % 2 == 1:  # 굵은 글씨
                                        run.bold = True
                else:
                    # 코드 블록
                    p = doc.add_paragraph(part, style='List Bullet')
                    p_format = p.paragraph_format
                    p_format.left_indent = Inches(0.5)
                    run = p.runs[0]
                    run.font.name = 'Consolas'
                    run.font.size = Pt(10)

        # 섹션 간 여백
        if level <= 2:
            doc.add_paragraph()

    # 저장
    doc.save(output_file)
    print(f"✓ DOCX 저장 완료: {output_file}")


def markdown_to_pptx(md_file, output_file):
    """Markdown → PPTX 변환"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    print(f"\n📊 PPTX 변환 시작: {md_file}")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 매뉴얼 파싱
    sections = parse_markdown(md_file)

    # 표지 슬라이드
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "GARAMe Manager 사용자 매뉴얼"
    subtitle.text = "v1.9.1\n\n산업안전 모니터링 시스템"

    # 본문 슬라이드
    current_slide = None
    bullet_level = 0

    for section in sections:
        level = section['level']
        title_text = section['title']
        content = section['content'].strip()

        # 목차나 메타정보 스킵
        if '목차' in title_text or 'Copyright' in title_text:
            continue

        # Level 1, 2는 새 슬라이드
        if level <= 2:
            # 제목 슬라이드
            blank_slide_layout = prs.slide_layouts[5]  # 빈 슬라이드
            slide = prs.slides.add_slide(blank_slide_layout)

            # 제목 텍스트박스
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(9)
            height = Inches(1)

            title_box = slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_box.text_frame
            title_frame.text = title_text

            p = title_frame.paragraphs[0]
            p.font.size = Pt(32) if level == 1 else Pt(28)
            p.font.bold = True
            p.font.color.rgb = RGBColor(44, 62, 80)
            p.alignment = PP_ALIGN.CENTER

            # 내용 텍스트박스
            if content:
                left = Inches(1)
                top = Inches(2)
                width = Inches(8)
                height = Inches(5)

                content_box = slide.shapes.add_textbox(left, top, width, height)
                content_frame = content_box.text_frame
                content_frame.word_wrap = True

                # 내용 추가 (간단하게 처리)
                lines = content.split('\n')
                for i, line in enumerate(lines[:10]):  # 최대 10줄
                    if line.strip():
                        if i == 0:
                            content_frame.text = line.strip()
                            p = content_frame.paragraphs[0]
                        else:
                            p = content_frame.add_paragraph()
                            p.text = line.strip()

                        p.font.size = Pt(14)
                        p.level = 0

                        # 리스트 처리
                        if re.match(r'^[-*]\s+', line):
                            p.text = line.strip()[2:]
                            p.level = 1

        # Level 3, 4는 내용만 추가 (이전 슬라이드에)
        elif level >= 3 and content and current_slide:
            try:
                # 기존 슬라이드에 내용 추가
                pass
            except:
                pass

        current_slide = slide

    # 저장
    prs.save(output_file)
    print(f"✓ PPTX 저장 완료: {output_file}")


def main():
    """메인 실행 함수"""
    import sys

    print("=" * 60)
    print("GARAMe Manager 매뉴얼 변환 스크립트")
    print("=" * 60)

    # 라이브러리 설치
    try:
        install_dependencies()
    except Exception as e:
        print(f"\n❌ 라이브러리 설치 실패: {e}")
        print("\n수동 설치 방법:")
        print("  pip3 install python-docx python-pptx markdown")
        return 1

    # 입력 파일
    md_file = Path(__file__).parent / "GARAMe_MANAGER_사용자매뉴얼_v1.9.1.md"

    if not md_file.exists():
        print(f"\n❌ 매뉴얼 파일을 찾을 수 없습니다: {md_file}")
        return 1

    # 출력 파일
    docx_file = md_file.with_suffix('.docx')
    pptx_file = md_file.with_suffix('.pptx')

    print(f"\n📝 입력 파일: {md_file.name}")
    print(f"📄 출력 파일 (DOCX): {docx_file.name}")
    print(f"📊 출력 파일 (PPTX): {pptx_file.name}")

    # 변환 수행
    try:
        # DOCX 변환
        markdown_to_docx(str(md_file), str(docx_file))

        # PPTX 변환
        markdown_to_pptx(str(md_file), str(pptx_file))

        print("\n" + "=" * 60)
        print("✅ 변환 완료!")
        print("=" * 60)
        print(f"\n생성된 파일:")
        print(f"  📄 {docx_file}")
        print(f"  📊 {pptx_file}")

        return 0

    except Exception as e:
        import traceback
        print(f"\n❌ 변환 실패: {e}")
        print(traceback.format_exc())
        return 1


if __name__ == "__main__":
    import sys
    sys.exit(main())
